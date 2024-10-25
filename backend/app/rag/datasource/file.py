import logging
import docx
import pptx
import openpyxl
from pydantic import BaseModel
from typing import Generator, IO
from pypdf import PdfReader

from app.models import Document, Upload
from app.file_storage import default_file_storage
from app.types import MimeTypes
from .base import BaseDataSource

logger = logging.getLogger(__name__)


class FileConfig(BaseModel):
    file_id: int


class FileDataSource(BaseDataSource):
    def validate_config(self):
        if not isinstance(self.config, list):
            raise ValueError("config must be a list")
        for f_config in self.config:
            FileConfig.model_validate(f_config)

    def load_documents(self) -> Generator[Document, None, None]:
        for f_config in self.config:
            upload_id = f_config["file_id"]
            upload = self.session.get(Upload, upload_id)
            if upload is None:
                logger.error(f"Upload with id {upload_id} not found")
                continue

            with default_file_storage.open(upload.path) as f:
                if upload.mime_type == MimeTypes.PDF:
                    content = extract_text_from_pdf(f)
                    mime_type = MimeTypes.PLAIN_TXT
                elif upload.mime_type == MimeTypes.DOCX:
                    content = extract_text_from_docx(f)
                    mime_type = MimeTypes.PLAIN_TXT
                elif upload.mime_type == MimeTypes.PPTX:
                    content = extract_text_from_pptx(f)
                    mime_type = MimeTypes.PLAIN_TXT
                elif upload.mime_type == MimeTypes.XLSX:
                    content = extract_text_from_xlsx(f)
                    mime_type = MimeTypes.PLAIN_TXT
                else:
                    content = f.read()
                    mime_type = upload.mime_type
            document = Document(
                name=upload.name,
                hash=hash(content),
                content=content,
                mime_type=mime_type,
                data_source_id=self.data_source_id,
                user_id=self.user_id,
                source_uri=upload.path,
                last_modified_at=upload.created_at,
            )
            yield document


def extract_text_from_pdf(file: IO) -> str:
    reader = PdfReader(file)
    return "\n\n".join([page.extract_text() for page in reader.pages])


def extract_text_from_docx(file: IO) -> str:
    document = docx.Document(file)
    full_text = []
    for paragraph in document.paragraphs:
        full_text.append(paragraph.text)
    return "\n\n".join(full_text)


def extract_text_from_pptx(file: IO) -> str:
    presentation = pptx.Presentation(file)
    full_text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text.append(shape.text)
    return "\n\n".join(full_text)


def extract_text_from_xlsx(file: IO) -> str:
    wb = openpyxl.load_workbook(file)
    full_text = []
    for sheet in wb.worksheets:
        full_text.append(f"Sheet: {sheet.title}")
        sheet_string = "\n".join(
            ",".join(map(str, row))
            for row in sheet.iter_rows(values_only=True)
        )
        full_text.append(sheet_string)
    return "\n\n".join(full_text)
